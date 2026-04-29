#!/usr/bin/env python3
"""阶段1：文档文字提取（独立快速运行）"""
import os, sys, json, time, traceback
from pathlib import Path

sys.stdout.reconfigure(encoding='utf-8', errors='replace')
sys.stderr.reconfigure(encoding='utf-8', errors='replace')

OUTPUT_ROOT = r"D:\东玄学习\转录文本"
DOC_DIRS = [
    r"D:\东玄学习\韩少嵛真神论命真神论命法\配套资料",
    r"D:\东玄学习\林白亲传弟子班课\命理风水 韩元茗稿",
    r"D:\东玄学习\林白亲传弟子班课\文档资料",
]
VIDEO_DIRS = [
    r"D:\东玄学习\韩元茗60天命名解码营40集\白泽易经大讲堂 姓名学--名字识人到取名旺运10集",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗《五行姓名学经典课程》（初级班 ）专栏 共三讲",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗60天命名解码营40集",
    r"D:\东玄学习\韩少嵛真神论命真神论命法",
]
PROGRESS_FILE = os.path.join(OUTPUT_ROOT, "_progress_v4.json")
LOG_FILE = os.path.join(OUTPUT_ROOT, "_transcribe_v4.log")
os.makedirs(OUTPUT_ROOT, exist_ok=True)

def log(msg):
    ts = time.strftime('%H:%M:%S')
    line = f"[{ts}] {msg}"
    print(line)
    with open(LOG_FILE, 'a', encoding='utf-8') as f:
        f.write(line + "\n")

def load_progress():
    if os.path.exists(PROGRESS_FILE):
        try:
            with open(PROGRESS_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except: pass
    return {"docs": {}, "videos": {}}

def save_progress(p):
    with open(PROGRESS_FILE, 'w', encoding='utf-8') as f:
        json.dump(p, f, ensure_ascii=False, indent=2)

def find_source_root_v2(filepath):
    all_dirs = VIDEO_DIRS + DOC_DIRS
    best_root, best_len = None, 0
    for src in all_dirs:
        try:
            rel = os.path.relpath(filepath, src)
            if not rel.startswith('..') and not os.path.isabs(rel):
                if len(src) > best_len:
                    best_root = src
                    best_len = len(src)
        except ValueError:
            continue
    return best_root

def get_output_path(source_path, source_root):
    rel = os.path.relpath(source_path, source_root)
    name = os.path.splitext(rel)[0]
    return os.path.normpath(os.path.join(OUTPUT_ROOT, name + '.txt'))

def is_done(output_txt, min_size=200):
    return os.path.exists(output_txt) and os.path.getsize(output_txt) >= min_size

# === 提取函数 ===
def extract_pdf_text(path):
    import pdfplumber
    texts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            t = page.extract_text()
            if t:
                texts.append(f"--- 第{i+1}页 ---\n{t}")
    return "\n\n".join(texts)

def extract_pptx_text(path):
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for i, slide in enumerate(prs.slides):
        st = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        st.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    rt = " | ".join(c.text.strip() for c in row.cells)
                    if rt.strip(" |"):
                        st.append(rt)
        if st:
            texts.append(f"--- 幻灯片 {i+1} ---\n" + "\n".join(st))
    return "\n\n".join(texts)

def extract_docx_text(path):
    from docx import Document
    doc = Document(path)
    texts = []
    for para in doc.paragraphs:
        txt = para.text.strip()
        if txt:
            texts.append(txt)
    for table in doc.tables:
        for row in table.rows:
            rt = " | ".join(c.text.strip() for c in row.cells)
            if rt.strip(" |"):
                texts.append(rt)
    return "\n".join(texts)

def extract_xlsx_text(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    texts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        st = [f"--- 工作表: {sheet} ---"]
        for row in ws.iter_rows(values_only=True):
            rt = " | ".join(str(c) if c else "" for c in row)
            if rt.strip(" |"):
                st.append(rt)
        texts.append("\n".join(st))
    wb.close()
    return "\n\n".join(texts)

def process_doc(fp, progress):
    sr = find_source_root_v2(fp) or os.path.dirname(fp)
    out = get_output_path(fp, sr)
    
    if is_done(out):
        progress["docs"][fp] = {"status": "skip_done", "output": out}
        return True, "skip"
    
    os.makedirs(os.path.dirname(out), exist_ok=True)
    ext = os.path.splitext(fp)[1].lower()
    
    try:
        if ext == ".pdf":
            text = extract_pdf_text(fp)
        elif ext == ".pptx":
            text = extract_pptx_text(fp)
        elif ext == ".docx":
            text = extract_docx_text(fp)
        elif ext == ".xlsx":
            text = extract_xlsx_text(fp)
        else:
            progress["docs"][fp] = {"status": "skip_ext", "output": out}
            return True, "skip_ext"
        
        if not text.strip():
            text = "(提取到的文字为空)"
        
        with open(out, 'w', encoding='utf-8') as f:
            f.write(f"# 文档文字提取\n# 源文件: {os.path.basename(fp)}\n")
            f.write(f"# 格式: {ext}\n")
            f.write(f"# 时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            f.write(text)
        
        size = os.path.getsize(out)
        progress["docs"][fp] = {"status": "done", "output": out, "size": size, "chars": len(text)}
        save_progress(progress)
        return True, "done"
    except Exception as e:
        err = str(e)
        progress["docs"][fp] = {"status": "fail", "error": err}
        save_progress(progress)
        return False, f"fail: {err[:80]}"

def scan_files(dirs, extensions):
    files = []
    for d in dirs:
        if not os.path.exists(d):
            log(f"[WARN] Dir not found: {d}")
            continue
        for root, _, fns in os.walk(d):
            for fn in fns:
                if os.path.splitext(fn)[1].lower() in extensions:
                    files.append(os.path.join(root, fn))
    return sorted(files)

def main():
    progress = load_progress()
    log("=" * 50)
    log("PHASE 1: 文档文字提取")
    log("=" * 50)
    
    doc_files = scan_files(DOC_DIRS, {".pdf", ".pptx", ".docx", ".xlsx"})
    log(f"扫描到 {len(doc_files)} 个文档")
    
    done = fail = skip = 0
    for i, fp in enumerate(doc_files):
        status = progress["docs"].get(fp, {}).get("status", "")
        if status == "done":
            skip += 1
            continue
        
        log(f"[{i+1}/{len(doc_files)}] {os.path.basename(fp)}")
        ok, result = process_doc(fp, progress)
        if ok:
            if result == "skip":
                skip += 1
            else:
                done += 1
                size = progress["docs"][fp].get("size", "?")
                chars = progress["docs"][fp].get("chars", "?")
                log(f"  -> OK ({size} bytes, {chars} chars)")
        else:
            fail += 1
            log(f"  -> FAIL: {result}")
    
    log("")
    log(f"DOCS DONE: ok={done} fail={fail} skip={skip}")
    
    # 统计输出文件
    txt_count = len([f for f in os.listdir(OUTPUT_ROOT) if f.endswith('.txt')])
    log(f"输出目录现有 .txt 文件: {txt_count}")
    log(f"进度保存: {PROGRESS_FILE}")
    save_progress(progress)

if __name__ == "__main__":
    main()
