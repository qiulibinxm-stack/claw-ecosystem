#!/usr/bin/env python3
"""
批量转录脚本 - 视频转文字 + 文档提取文字
支持：MP4→音频→Whisper转录→TXT, PDF/PPTX/DOCX/XLSX→TXT
"""

import os
import sys
import json
import subprocess
import time
import traceback
from pathlib import Path

# === 配置 ===
FFMPEG_PATH = r"C:\Program Files\EVCapture\ffmpeg.exe"
WHISPER_MODEL = "medium"  # tiny/base/small/medium/large/turbo
WHISPER_LANGUAGE = "zh"   # 中文
OUTPUT_ROOT = r"D:\东玄学习\转录文本"

# 源目录配置
VIDEO_DIRS = [
    r"D:\东玄学习\韩元茗60天命名解码营40集\白泽易经大讲堂 姓名学--名字识人到取名旺运10集",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗《五行姓名学经典课程》（初级班 ）专栏 共三讲",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗《五行姓名学经典课程》（高级班 ）专栏 共九讲",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗60天命名解码营40集",
    r"D:\东玄学习\韩少嵛真神论命真神论命法",
]

DOC_DIRS = [
    r"D:\东玄学习\韩少嵛真神论命真神论命法\配套资料",
    r"D:\东玄学习\林白亲传弟子班课\命理风水 韩元茗稿",
    r"D:\东玄学习\林白亲传弟子班课\文档资料",
]

# 进度日志
PROGRESS_FILE = os.path.join(OUTPUT_ROOT, "_progress.json")

def load_progress():
    """加载处理进度"""
    if os.path.exists(PROGRESS_FILE):
        with open(PROGRESS_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {"completed": [], "failed": []}

def save_progress(progress):
    """保存处理进度"""
    os.makedirs(os.path.dirname(PROGRESS_FILE), exist_ok=True)
    with open(PROGRESS_FILE, 'w', encoding='utf-8') as f:
        json.dump(progress, f, ensure_ascii=False, indent=2)

def get_output_path(source_path, source_root):
    """根据源文件路径计算输出TXT路径"""
    rel = os.path.relpath(source_path, source_root)
    # 去掉原始扩展名，加.txt
    rel_txt = os.path.splitext(rel)[0] + ".txt"
    return os.path.join(OUTPUT_ROOT, rel_txt)

def find_source_root(filepath, dirs):
    """找到文件属于哪个源目录"""
    for d in dirs:
        try:
            os.path.relpath(filepath, d)
            return d
        except ValueError:
            continue
    return None

def extract_audio(video_path, audio_path):
    """用ffmpeg从视频中提取音频（16kHz mono WAV，Whisper最佳输入格式）"""
    cmd = [
        FFMPEG_PATH, "-y", "-i", video_path,
        "-vn",                    # 不要视频
        "-acodec", "pcm_s16le",   # 16-bit PCM
        "-ar", "16000",           # 16kHz采样率
        "-ac", "1",               # 单声道
        audio_path
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg failed: {result.stderr[-500:]}")
    return audio_path

def transcribe_audio(audio_path, output_txt, model=WHISPER_MODEL, language=WHISPER_LANGUAGE):
    """用whisper转录音频"""
    import whisper
    
    print(f"  加载Whisper模型: {model}...")
    model_obj = whisper.load_model(model)
    
    print(f"  转录中（语言={language}）...")
    result = model_obj.transcribe(audio_path, language=language, verbose=False)
    
    text = result["text"].strip()
    
    # 也保存分段信息
    segments = result.get("segments", [])
    
    with open(output_txt, 'w', encoding='utf-8') as f:
        f.write(f"# 转录文本\n")
        f.write(f"# 模型: {model}\n")
        f.write(f"# 时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
        
        # 全文
        f.write("=" * 60 + "\n")
        f.write("全文\n")
        f.write("=" * 60 + "\n\n")
        f.write(text + "\n\n")
        
        # 带时间戳的分段
        f.write("=" * 60 + "\n")
        f.write("分段文本（带时间戳）\n")
        f.write("=" * 60 + "\n\n")
        for seg in segments:
            start = format_timestamp(seg["start"])
            end = format_timestamp(seg["end"])
            f.write(f"[{start} → {end}] {seg['text'].strip()}\n")
    
    return output_txt

def format_timestamp(seconds):
    """格式化时间戳"""
    h = int(seconds // 3600)
    m = int((seconds % 3600) // 60)
    s = int(seconds % 60)
    return f"{h:02d}:{m:02d}:{s:02d}"

def process_video(video_path, progress):
    """处理单个视频文件"""
    all_dirs = VIDEO_DIRS + DOC_DIRS
    source_root = find_source_root(video_path, all_dirs)
    if not source_root:
        source_root = os.path.dirname(video_path)
    
    output_txt = get_output_path(video_path, source_root)
    
    # 跳过已完成的
    if video_path in progress["completed"]:
        print(f"  [跳过] 已完成: {os.path.basename(video_path)}")
        return True
    
    # 如果输出文件已存在
    if os.path.exists(output_txt) and os.path.getsize(output_txt) > 100:
        print(f"  [跳过] 输出已存在: {os.path.basename(output_txt)}")
        progress["completed"].append(video_path)
        save_progress(progress)
        return True
    
    os.makedirs(os.path.dirname(output_txt), exist_ok=True)
    
    # 临时音频文件
    temp_audio = output_txt.replace(".txt", ".wav")
    
    try:
        # Step 1: 提取音频
        print(f"  [1/2] 提取音频...")
        extract_audio(video_path, temp_audio)
        
        # Step 2: Whisper转录
        print(f"  [2/2] Whisper转录...")
        transcribe_audio(temp_audio, output_txt)
        
        # 清理临时音频
        if os.path.exists(temp_audio):
            os.remove(temp_audio)
        
        progress["completed"].append(video_path)
        save_progress(progress)
        print(f"  ✅ 完成: {os.path.basename(output_txt)}")
        return True
        
    except Exception as e:
        print(f"  ❌ 失败: {e}")
        traceback.print_exc()
        progress["failed"].append({"file": video_path, "error": str(e)})
        save_progress(progress)
        # 清理临时文件
        if os.path.exists(temp_audio):
            try: os.remove(temp_audio)
            except: pass
        return False

def extract_pdf_text(pdf_path):
    """提取PDF文字"""
    import pdfplumber
    texts = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                texts.append(f"--- 第{i+1}页 ---\n{text}")
    return "\n\n".join(texts)

def extract_pptx_text(pptx_path):
    """提取PPTX文字"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    texts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            # 表格
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_texts.append(row_text)
        if slide_texts:
            texts.append(f"--- 幻灯片 {i+1} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)

def extract_docx_text(docx_path):
    """提取DOCX文字"""
    from docx import Document
    doc = Document(docx_path)
    texts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            texts.append(text)
    # 表格
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                texts.append(row_text)
    return "\n".join(texts)

def extract_xlsx_text(xlsx_path):
    """提取XLSX文字"""
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
    texts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_texts = [f"--- 工作表: {sheet_name} ---"]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip(" |"):
                sheet_texts.append(row_text)
        texts.append("\n".join(sheet_texts))
    wb.close()
    return "\n\n".join(texts)

def process_document(doc_path, progress):
    """处理单个文档文件"""
    all_dirs = VIDEO_DIRS + DOC_DIRS
    source_root = find_source_root(doc_path, all_dirs)
    if not source_root:
        source_root = os.path.dirname(doc_path)
    
    output_txt = get_output_path(doc_path, source_root)
    
    # 跳过已完成的
    if doc_path in progress["completed"]:
        print(f"  [跳过] 已完成: {os.path.basename(doc_path)}")
        return True
    
    if os.path.exists(output_txt) and os.path.getsize(output_txt) > 50:
        print(f"  [跳过] 输出已存在: {os.path.basename(output_txt)}")
        progress["completed"].append(doc_path)
        save_progress(progress)
        return True
    
    os.makedirs(os.path.dirname(output_txt), exist_ok=True)
    
    ext = os.path.splitext(doc_path)[1].lower()
    
    try:
        print(f"  提取文字...")
        if ext == ".pdf":
            text = extract_pdf_text(doc_path)
        elif ext == ".pptx":
            text = extract_pptx_text(doc_path)
        elif ext == ".docx":
            text = extract_docx_text(doc_path)
        elif ext == ".xlsx":
            text = extract_xlsx_text(doc_path)
        else:
            print(f"  [跳过] 不支持的格式: {ext}")
            return True
        
        if not text.strip():
            text = "（提取到的文字为空）"
        
        with open(output_txt, 'w', encoding='utf-8') as f:
            f.write(f"# 文档文字提取\n")
            f.write(f"# 源文件: {os.path.basename(doc_path)}\n")
            f.write(f"# 时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            f.write(text)
        
        progress["completed"].append(doc_path)
        save_progress(progress)
        print(f"  ✅ 完成: {os.path.basename(output_txt)}")
        return True
        
    except Exception as e:
        print(f"  ❌ 失败: {e}")
        traceback.print_exc()
        progress["failed"].append({"file": doc_path, "error": str(e)})
        save_progress(progress)
        return False

def scan_files(dirs, extensions):
    """扫描目录中的文件"""
    files = []
    for d in dirs:
        if not os.path.exists(d):
            print(f"  [警告] 目录不存在: {d}")
            continue
        for root, _, filenames in os.walk(d):
            for fn in filenames:
                if os.path.splitext(fn)[1].lower() in extensions:
                    files.append(os.path.join(root, fn))
    return sorted(files)

def main():
    print("=" * 70)
    print("  批量转录 & 文档文字提取工具")
    print("=" * 70)
    print(f"  Whisper模型: {WHISPER_MODEL}")
    print(f"  语言: {WHISPER_LANGUAGE}")
    print(f"  输出目录: {OUTPUT_ROOT}")
    print()
    
    # 确保ffmpeg存在
    if not os.path.exists(FFMPEG_PATH):
        print(f"❌ ffmpeg不存在: {FFMPEG_PATH}")
        sys.exit(1)
    
    # 设置环境变量让whisper找到ffmpeg
    os.environ["PATH"] = os.path.dirname(FFMPEG_PATH) + os.pathsep + os.environ.get("PATH", "")
    
    progress = load_progress()
    
    # === 1. 扫描视频文件 ===
    print("📂 扫描视频文件...")
    video_files = scan_files(VIDEO_DIRS, {".mp4", ".avi", ".mkv", ".mov", ".flv"})
    print(f"  找到 {len(video_files)} 个视频文件")
    
    # === 2. 扫描文档文件 ===
    print("📂 扫描文档文件...")
    doc_files = scan_files(DOC_DIRS, {".pdf", ".pptx", ".docx", ".xlsx"})
    print(f"  找到 {len(doc_files)} 个文档文件")
    
    # === 3. 先处理文档（快） ===
    print(f"\n{'='*70}")
    print(f"  阶段1：文档文字提取 ({len(doc_files)} 个)")
    print(f"{'='*70}")
    
    doc_done = 0
    doc_fail = 0
    for i, doc_path in enumerate(doc_files):
        print(f"\n[{i+1}/{len(doc_files)}] {os.path.basename(doc_path)}")
        if process_document(doc_path, progress):
            doc_done += 1
        else:
            doc_fail += 1
    
    print(f"\n📊 文档处理完成: ✅{doc_done} ❌{doc_fail}")
    
    # === 4. 处理视频（慢） ===
    print(f"\n{'='*70}")
    print(f"  阶段2：视频转录 ({len(video_files)} 个)")
    print(f"  ⏱️ 预计耗时较长，请耐心等待...")
    print(f"{'='*70}")
    
    vid_done = 0
    vid_fail = 0
    start_time = time.time()
    
    for i, video_path in enumerate(video_files):
        elapsed = time.time() - start_time
        if vid_done > 0:
            avg_time = elapsed / vid_done
            remaining = avg_time * (len(video_files) - i)
            eta_str = f"预计剩余: {remaining/60:.0f}分钟"
        else:
            eta_str = ""
        
        size_mb = os.path.getsize(video_path) / (1024*1024)
        print(f"\n[{i+1}/{len(video_files)}] {os.path.basename(video_path)} ({size_mb:.0f}MB) {eta_str}")
        
        if process_video(video_path, progress):
            vid_done += 1
        else:
            vid_fail += 1
    
    total_time = time.time() - start_time
    
    # === 5. 汇总 ===
    print(f"\n{'='*70}")
    print(f"  🏁 全部完成！")
    print(f"{'='*70}")
    print(f"  文档: ✅{doc_done} ❌{doc_fail}")
    print(f"  视频: ✅{vid_done} ❌{vid_fail}")
    print(f"  总耗时: {total_time/60:.1f}分钟")
    print(f"  输出目录: {OUTPUT_ROOT}")
    
    if progress["failed"]:
        print(f"\n❌ 失败列表:")
        for item in progress["failed"]:
            print(f"  - {item['file']}: {item['error']}")

if __name__ == "__main__":
    main()
