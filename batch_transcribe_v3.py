#!/usr/bin/env python3
"""
批量转录脚本 v3 - Gemini API视频转录 + 文档文字提取
修复了v2的进度bug：文件写出后才标记完成
"""
import os, sys, json, subprocess, time, base64, traceback, urllib.request, urllib.error
from pathlib import Path

sys.stdout.reconfigure(encoding='utf-8', errors='replace')
sys.stderr.reconfigure(encoding='utf-8', errors='replace')

# === 配置 ===
FFMPEG_PATH = r"C:\Program Files\EVCapture\ffmpeg.exe"
OUTPUT_ROOT = r"D:\东玄学习\转录文本"
GEMINI_MODEL = "gemini-2.5-flash"
GEMINI_MAX_AUDIO_BYTES = 18 * 1024 * 1024

# 源目录
VIDEO_DIRS = [
    r"D:\东玄学习\韩元茗60天命名解码营40集\白泽易经大讲堂 姓名学--名字识人到取名旺运10集",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗《五行姓名学经典课程》（初级班 ）专栏 共三讲",
    # 高级班目录不存在，跳过
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗60天命名解码营40集",
    r"D:\东玄学习\韩少嵛真神论命真神论命法",
]
DOC_DIRS = [
    r"D:\东玄学习\韩少嵛真神论命真神论命法\配套资料",
    r"D:\东玄学习\林白亲传弟子班课\命理风水 韩元茗稿",
    r"D:\东玄学习\林白亲传弟子班课\文档资料",
]

PROGRESS_FILE = os.path.join(OUTPUT_ROOT, "_progress.json")
LOG_FILE = os.path.join(OUTPUT_ROOT, "_transcribe.log")
os.makedirs(OUTPUT_ROOT, exist_ok=True)

def log(msg):
    ts = time.strftime('%Y-%m-%d %H:%M:%S')
    line = f"[{ts}] {msg}"
    print(line)
    with open(LOG_FILE, 'a', encoding='utf-8') as f:
        f.write(line + "\n")

def load_progress():
    if os.path.exists(PROGRESS_FILE):
        try:
            with open(PROGRESS_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            pass
    return {"completed": [], "failed": []}

def save_progress(progress):
    with open(PROGRESS_FILE, 'w', encoding='utf-8') as f:
        json.dump(progress, f, ensure_ascii=False, indent=2)

def load_gemini_key():
    with open(r"C:\Users\Administrator\.qclaw\openclaw.json", 'r', encoding='utf-8') as f:
        config = json.load(f)
    return config['models']['providers']['gemini']['apiKey']

def find_source_root(filepath, dirs):
    for d in dirs:
        try:
            os.path.relpath(filepath, d)
            return d
        except ValueError:
            continue
    return os.path.dirname(filepath)

def get_output_path(source_path, source_root):
    rel = os.path.relpath(source_path, source_root)
    rel_txt = os.path.splitext(rel)[0] + ".txt"
    return os.path.join(OUTPUT_ROOT, rel_txt)

def is_actually_done(output_txt):
    """检查是否真正完成（文件存在且大小合理）"""
    return os.path.exists(output_txt) and os.path.getsize(output_txt) > 100

# === 视频处理 ===
def get_video_duration(video_path):
    cmd = [FFMPEG_PATH, "-i", video_path]
    result = subprocess.run(cmd, capture_output=True, timeout=30)
    stderr = result.stderr.decode('utf-8', errors='replace') if result.stderr else ''
    for line in stderr.split('\n'):
        if 'Duration' in line:
            try:
                dur = line.split('Duration:')[1].split(',')[0].strip()
                h, m, s = dur.split(':')
                return int(float(h)*3600 + float(m)*60 + float(s))
            except:
                pass
    return 0

def extract_audio(video_path, audio_path, start=0, duration=None):
    cmd = [FFMPEG_PATH, "-y", "-i", video_path]
    if start > 0:
        cmd.extend(["-ss", str(start)])
    if duration:
        cmd.extend(["-t", str(duration)])
    cmd.extend(["-vn", "-acodec", "libmp3lame", "-ab", "64k", "-ar", "16000", "-ac", "1", audio_path])
    r = subprocess.run(cmd, capture_output=True, timeout=600)
    if r.returncode != 0:
        err = r.stderr.decode('utf-8', errors='replace')[-200:] if r.stderr else 'unknown'
        raise RuntimeError(f"ffmpeg failed: {err}")

def gemini_transcribe(audio_path, chunk_info=""):
    with open(audio_path, 'rb') as f:
        audio_data = f.read()
    audio_b64 = base64.b64encode(audio_data).decode('utf-8')
    
    prompt = f"""请将这段音频完整转录为文字。要求：
1. 完整转录所有说的内容，不要遗漏
2. 保留说话人的语气和表达方式
3. 如有专业术语，保留原词
4. 按自然段落分段
{chunk_info}"""
    
    payload = {
        "contents": [{"parts": [{"text": prompt}, {"inline_data": {"mime_type": "audio/mp3", "data": audio_b64}}]}],
        "generationConfig": {"temperature": 0.1, "maxOutputTokens": 65536}
    }
    
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent?key={GEMINI_API_KEY}"
    req = urllib.request.Request(url, data=json.dumps(payload).encode('utf-8'),
                                headers={"Content-Type": "application/json"}, method="POST")
    
    for attempt in range(3):
        try:
            with urllib.request.urlopen(req, timeout=180) as resp:
                result = json.loads(resp.read().decode('utf-8'))
                return result['candidates'][0]['content']['parts'][0]['text'].strip()
        except urllib.error.HTTPError as e:
            body = e.read().decode('utf-8', errors='replace') if e.fp else ''
            if e.code == 429:
                log(f"  Rate limit, waiting {30*(attempt+1)}s...")
                time.sleep(30*(attempt+1))
            elif e.code == 503:
                log(f"  Service overloaded, waiting {60*(attempt+1)}s...")
                time.sleep(60*(attempt+1))
            else:
                raise RuntimeError(f"HTTP {e.code}: {body[:200]}")
        except Exception as e:
            if attempt < 2:
                log(f"  Retry {attempt+1}/3: {e}")
                time.sleep(10)
                continue
            raise
    raise RuntimeError("Gemini API retries exhausted")

def process_video(video_path, progress):
    all_dirs = VIDEO_DIRS + DOC_DIRS
    source_root = find_source_root(video_path, all_dirs)
    output_txt = get_output_path(video_path, source_root)
    
    # 真实完成检查（用output_txt判断）
    if is_actually_done(output_txt):
        if output_txt not in progress["completed"]:
            progress["completed"].append(output_txt)
            save_progress(progress)
        return True
    
    os.makedirs(os.path.dirname(output_txt), exist_ok=True)
    duration = get_video_duration(video_path)
    log(f"  Duration: {duration//60}m{duration%60}s")
    
    max_chunk = 25 * 60  # 25min per chunk
    chunks = []
    
    if duration <= max_chunk:
        temp_audio = output_txt.replace(".txt", "_t.mp3")
        try:
            log(f"  [1/2] Extract audio...")
            extract_audio(video_path, temp_audio)
            sz = os.path.getsize(temp_audio) / (1024*1024)
            log(f"  Audio: {sz:.1f}MB")
            
            if sz > 18:
                log(f"  Audio too large, splitting...")
                os.remove(temp_audio)
                return process_video_chunked(video_path, output_txt, duration, 10*60, progress)
            
            log(f"  [2/2] Gemini transcribe...")
            text = gemini_transcribe(temp_audio)
        finally:
            if os.path.exists(temp_audio):
                try: os.remove(temp_audio)
                except: pass
    else:
        return process_video_chunked(video_path, output_txt, duration, max_chunk, progress)
    
    # 写文件后才标记完成
    with open(output_txt, 'w', encoding='utf-8') as f:
        f.write(f"# 转录文本\n# 源文件: {os.path.basename(video_path)}\n")
        f.write(f"# 时长: {duration//60}分{duration%60}秒\n")
        f.write(f"# 转录方式: Gemini API\n")
        f.write(f"# 时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
        f.write(text)
    
    if output_txt not in progress["completed"]:
        progress["completed"].append(video_path)
        save_progress(progress)
    return True

def process_video_chunked(video_path, output_txt, duration, chunk_dur, progress):
    num_chunks = (duration + chunk_dur - 1) // chunk_dur
    chunks_text = []
    
    for i in range(num_chunks):
        start = i * chunk_dur
        dur = min(chunk_dur, duration - start)
        temp_audio = output_txt.replace(".txt", f"_c{i+1}.mp3")
        
        try:
            log(f"  Chunk {i+1}/{num_chunks}: {start//60}m-{dur//60}m")
            extract_audio(video_path, temp_audio, start, dur)
            info = f"(chunk {i+1}/{num_chunks}, ~{start//60}min to {(start+dur)//60}min)"
            text = gemini_transcribe(temp_audio, info)
            chunks_text.append(text)
        except Exception as e:
            log(f"  Chunk {i+1} failed: {e}")
            chunks_text.append(f"[Chunk {i+1} failed]")
        finally:
            if os.path.exists(temp_audio):
                try: os.remove(temp_audio)
                except: pass
        
        if i < num_chunks - 1:
            time.sleep(5)
    
    full_text = "\n\n".join(chunks_text)
    with open(output_txt, 'w', encoding='utf-8') as f:
        f.write(f"# 转录文本\n# 源文件: {os.path.basename(video_path)}\n")
        f.write(f"# 时长: {duration//60}分{duration%60}秒\n")
        f.write(f"# 转录方式: Gemini API\n")
        f.write(f"# 分段: {num_chunks}\n")
        f.write(f"# 时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
        f.write(full_text)
    
    if output_txt not in progress["completed"]:
        progress["completed"].append(video_path)
        save_progress(progress)
    return True

# === 文档处理 ===
def extract_pdf_text(pdf_path):
    import pdfplumber
    texts = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            t = page.extract_text()
            if t:
                texts.append(f"--- 第{i+1}页 ---\n{t}")
    return "\n\n".join(texts)

def extract_pptx_text(pptx_path):
    from pptx import Presentation
    prs = Presentation(pptx_path)
    texts = []
    for i, slide in enumerate(prs.slides):
        st = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        st.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    rt = " | ".join(c.text.strip() for c in row.cells)
                    if rt.strip(" |"):
                        st.append(rt)
        if st:
            texts.append(f"--- 幻灯片 {i+1} ---\n" + "\n".join(st))
    return "\n\n".join(texts)

def extract_docx_text(docx_path):
    from docx import Document
    doc = Document(docx_path)
    texts = []
    for para in doc.paragraphs:
        txt = para.text.strip()
        if txt:
            texts.append(txt)
    for table in doc.tables:
        for row in table.rows:
            rt = " | ".join(c.text.strip() for c in row.cells)
            if rt.strip(" |"):
                texts.append(rt)
    return "\n".join(texts)

def extract_xlsx_text(xlsx_path):
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
    texts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        st = [f"--- 工作表: {sheet} ---"]
        for row in ws.iter_rows(values_only=True):
            rt = " | ".join(str(c) if c else "" for c in row)
            if rt.strip(" |"):
                st.append(rt)
        texts.append("\n".join(st))
    wb.close()
    return "\n\n".join(texts)

def process_document(doc_path, progress):
    all_dirs = VIDEO_DIRS + DOC_DIRS
    source_root = find_source_root(doc_path, all_dirs)
    output_txt = get_output_path(doc_path, source_root)
    
    # 用output_txt路径来判断是否完成（更可靠）
    if is_actually_done(output_txt):
        # 确保在progress里
        if output_txt not in progress["completed"]:
            progress["completed"].append(output_txt)
            save_progress(progress)
        return True
    
    os.makedirs(os.path.dirname(output_txt), exist_ok=True)
    ext = os.path.splitext(doc_path)[1].lower()
    
    try:
        if ext == ".pdf":
            text = extract_pdf_text(doc_path)
        elif ext == ".pptx":
            text = extract_pptx_text(doc_path)
        elif ext == ".docx":
            text = extract_docx_text(doc_path)
        elif ext == ".xlsx":
            text = extract_xlsx_text(doc_path)
        else:
            return True
        
        if not text.strip():
            text = "(提取到的文字为空)"
        
        with open(output_txt, 'w', encoding='utf-8') as f:
            f.write(f"# 文档文字提取\n# 源文件: {os.path.basename(doc_path)}\n")
            f.write(f"# 格式: {ext}\n")
            f.write(f"# 时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            f.write(text)
        
        if output_txt not in progress["completed"]:
            progress["completed"].append(doc_path)
            save_progress(progress)
        return True
    except Exception as e:
        log(f"  Doc failed: {e}")
        progress["failed"].append({"file": doc_path, "error": str(e)})
        save_progress(progress)
        return False

def scan_files(dirs, extensions):
    files = []
    for d in dirs:
        if not os.path.exists(d):
            log(f"  [WARN] Dir not found: {d}")
            continue
        for root, _, fns in os.walk(d):
            for fn in fns:
                if os.path.splitext(fn)[1].lower() in extensions:
                    files.append(os.path.join(root, fn))
    return sorted(files)

def main():
    global GEMINI_API_KEY
    
    print("=" * 60)
    print("  批量转录 v3 - Gemini API + 文档提取")
    print("=" * 60)
    
    GEMINI_API_KEY = load_gemini_key()
    if not GEMINI_API_KEY:
        print("ERROR: Gemini API key not found")
        sys.exit(1)
    print(f"  API Key: {GEMINI_API_KEY[:10]}...")
    print(f"  Model: {GEMINI_MODEL}")
    print(f"  Output: {OUTPUT_ROOT}")
    
    progress = load_progress()
    log("=" * 60)
    log("SCANNING FILES...")
    log("=" * 60)
    
    video_files = scan_files(VIDEO_DIRS, {".mp4", ".avi", ".mkv", ".mov", ".flv"})
    doc_files = scan_files(DOC_DIRS, {".pdf", ".pptx", ".docx", ".xlsx"})
    log(f"  Videos: {len(video_files)}")
    log(f"  Docs: {len(doc_files)}")
    
    # === PHASE 1: Docs (fast) ===
    log(f"\n{'='*60}")
    log(f"PHASE 1: DOCS ({len(doc_files)})")
    log(f"{'='*60}")
    
    doc_done = doc_fail = 0
    for i, fp in enumerate(doc_files):
        log(f"[{i+1}/{len(doc_files)}] {os.path.basename(fp)}")
        if process_document(fp, progress):
            doc_done += 1
        else:
            doc_fail += 1
    
    log(f"DOCS DONE: ok={doc_done} fail={doc_fail}")
    
    # === PHASE 2: Videos (slow) ===
    log(f"\n{'='*60}")
    log(f"PHASE 2: VIDEOS ({len(video_files)})")
    log(f"{'='*60}")
    
    vid_done = vid_fail = 0
    start = time.time()
    
    for i, vp in enumerate(video_files):
        elapsed = time.time() - start
        eta = f"ETA ~{elapsed/vid_done*(len(video_files)-i)/60:.0f}min" if vid_done > 0 else ""
        sz = os.path.getsize(vp) / (1024*1024)
        log(f"\n[{i+1}/{len(video_files)}] {os.path.basename(vp)} ({sz:.0f}MB) {eta}")
        
        if process_video(vp, progress):
            vid_done += 1
        else:
            vid_fail += 1
        
        time.sleep(3)
    
    total = time.time() - start
    log(f"\n{'='*60}")
    log(f"ALL DONE!")
    log(f"  Docs: ok={doc_done} fail={doc_fail}")
    log(f"  Videos: ok={vid_done} fail={vid_fail}")
    log(f"  Time: {total/60:.1f}min")
    log(f"  Output: {OUTPUT_ROOT}")
    log(f"{'='*60}")
    
    if progress["failed"]:
        log("FAILURES:")
        for item in progress["failed"]:
            log(f"  - {os.path.basename(item['file'])}: {item['error']}")

if __name__ == "__main__":
    main()
