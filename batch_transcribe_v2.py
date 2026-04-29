#!/usr/bin/env python3
"""
批量转录脚本 v2 - 用Gemini API做视频转录 + 本地文档文字提取
支持：MP4→音频→Gemini转录→TXT, PDF/PPTX/DOCX/XLSX→TXT
"""

import os
import sys
import json
import subprocess
import time
import base64
import traceback
import urllib.request
import urllib.error
from pathlib import Path

# Fix Windows GBK encoding
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
if sys.stderr.encoding != 'utf-8':
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')

# === 配置 ===
FFMPEG_PATH = r"C:\Program Files\EVCapture\ffmpeg.exe"
OUTPUT_ROOT = r"D:\东玄学习\转录文本"

# Gemini API配置
GEMINI_API_KEY = None  # 从openclaw.json读取
GEMINI_MODEL = "gemini-2.5-flash"  # 快速且支持音频
GEMINI_MAX_AUDIO_BYTES = 19 * 1024 * 1024  # 19MB安全限制(Gemini内联数据限制20MB)

# 源目录配置
VIDEO_DIRS = [
    r"D:\东玄学习\韩元茗60天命名解码营40集\白泽易经大讲堂 姓名学--名字识人到取名旺运10集",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗《五行姓名学经典课程》（初级班 ）专栏 共三讲",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗《五行姓名学经典课程》（高级班 ）专栏 共九讲",
    r"D:\东玄学习\韩元茗60天命名解码营40集\韩元茗60天命名解码营40集",
    r"D:\东玄学习\韩少嵛真神论命真神论命法",
]

DOC_DIRS = [
    r"D:\东玄学习\韩少嵛真神论命真神论命法\配套资料",
    r"D:\东玄学习\林白亲传弟子班课\命理风水 韩元茗稿",
    r"D:\东玄学习\林白亲传弟子班课\文档资料",
]

# 进度日志
PROGRESS_FILE = os.path.join(OUTPUT_ROOT, "_progress.json")
LOG_FILE = os.path.join(OUTPUT_ROOT, "_transcribe.log")

def log(msg):
    """写日志"""
    timestamp = time.strftime('%Y-%m-%d %H:%M:%S')
    line = f"[{timestamp}] {msg}"
    print(line)
    os.makedirs(os.path.dirname(LOG_FILE), exist_ok=True)
    with open(LOG_FILE, 'a', encoding='utf-8') as f:
        f.write(line + "\n")

def load_progress():
    if os.path.exists(PROGRESS_FILE):
        with open(PROGRESS_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {"completed": [], "failed": [], "stats": {}}

def save_progress(progress):
    os.makedirs(os.path.dirname(PROGRESS_FILE), exist_ok=True)
    with open(PROGRESS_FILE, 'w', encoding='utf-8') as f:
        json.dump(progress, f, ensure_ascii=False, indent=2)

def load_gemini_key():
    """从openclaw.json读取Gemini API key"""
    config_path = r"C:\Users\Administrator\.qclaw\openclaw.json"
    with open(config_path, 'r', encoding='utf-8') as f:
        config = json.load(f)
    return config.get('models', {}).get('providers', {}).get('gemini', {}).get('apiKey', '')

def get_output_path(source_path, source_root):
    rel = os.path.relpath(source_path, source_root)
    rel_txt = os.path.splitext(rel)[0] + ".txt"
    return os.path.join(OUTPUT_ROOT, rel_txt)

def find_source_root(filepath, dirs):
    for d in dirs:
        try:
            os.path.relpath(filepath, d)
            return d
        except ValueError:
            continue
    return os.path.dirname(filepath)

def extract_audio_segment(video_path, audio_path, start_sec=0, duration_sec=None):
    """用ffmpeg从视频中提取音频段（MP3格式，节省空间）"""
    cmd = [FFMPEG_PATH, "-y", "-i", video_path]
    if start_sec > 0:
        cmd.extend(["-ss", str(start_sec)])
    if duration_sec:
        cmd.extend(["-t", str(duration_sec)])
    cmd.extend([
        "-vn",
        "-acodec", "libmp3lame",
        "-ab", "64k",      # 64kbps 足够语音识别
        "-ar", "16000",     # 16kHz
        "-ac", "1",         # 单声道
        audio_path
    ])
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg failed: {result.stderr[-300:]}")
    return audio_path

def get_video_duration(video_path):
    """获取视频时长（秒）"""
    cmd = [FFMPEG_PATH, "-i", video_path, "-hide_banner"]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
    # 解析 Duration: 01:23:45.67
    for line in result.stderr.split('\n'):
        if 'Duration' in line:
            try:
                dur_str = line.split('Duration:')[1].split(',')[0].strip()
                h, m, s = dur_str.split(':')
                return int(float(h) * 3600 + float(m) * 60 + float(s))
            except:
                pass
    return 0

def gemini_transcribe(audio_path, chunk_info=""):
    """使用Gemini API转录音频"""
    with open(audio_path, 'rb') as f:
        audio_data = f.read()
    
    audio_b64 = base64.b64encode(audio_data).decode('utf-8')
    
    # 判断MIME类型
    if audio_path.endswith('.mp3'):
        mime_type = "audio/mp3"
    elif audio_path.endswith('.wav'):
        mime_type = "audio/wav"
    else:
        mime_type = "audio/mp3"
    
    prompt = f"""请将这段音频完整转录为文字。要求：
1. 完整转录所有说的内容，不要遗漏
2. 保留说话人的语气和表达方式
3. 如有专业术语，保留原词
4. 按自然段落分段
{chunk_info}"""

    payload = {
        "contents": [{
            "parts": [
                {"text": prompt},
                {"inline_data": {"mime_type": mime_type, "data": audio_b64}}
            ]
        }],
        "generationConfig": {
            "temperature": 0.1,
            "maxOutputTokens": 65536
        }
    }
    
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent?key={GEMINI_API_KEY}"
    
    req = urllib.request.Request(
        url,
        data=json.dumps(payload).encode('utf-8'),
        headers={"Content-Type": "application/json"},
        method="POST"
    )
    
    max_retries = 3
    for attempt in range(max_retries):
        try:
            with urllib.request.urlopen(req, timeout=120) as resp:
                result = json.loads(resp.read().decode('utf-8'))
                text = result['candidates'][0]['content']['parts'][0]['text']
                return text.strip()
        except urllib.error.HTTPError as e:
            error_body = e.read().decode('utf-8') if e.fp else ''
            if e.code == 429:  # Rate limit
                wait_time = 30 * (attempt + 1)
                log(f"  ⏳ 限流，等待{wait_time}秒...")
                time.sleep(wait_time)
                continue
            elif e.code == 503:  # Overloaded
                wait_time = 60 * (attempt + 1)
                log(f"  ⏳ 服务过载，等待{wait_time}秒...")
                time.sleep(wait_time)
                continue
            else:
                raise RuntimeError(f"Gemini API error {e.code}: {error_body[:300]}")
        except Exception as e:
            if attempt < max_retries - 1:
                log(f"  ⚠️ 重试 ({attempt+1}/{max_retries}): {e}")
                time.sleep(10)
                continue
            raise
    
    raise RuntimeError("Gemini API重试次数耗尽")

def process_video(video_path, progress):
    """处理单个视频文件"""
    all_dirs = VIDEO_DIRS + DOC_DIRS
    source_root = find_source_root(video_path, all_dirs)
    output_txt = get_output_path(video_path, source_root)
    
    # 跳过已完成的
    if video_path in progress["completed"]:
        return True
    if os.path.exists(output_txt) and os.path.getsize(output_txt) > 100:
        progress["completed"].append(video_path)
        save_progress(progress)
        return True
    
    os.makedirs(os.path.dirname(output_txt), exist_ok=True)
    
    # 获取视频时长
    duration = get_video_duration(video_path)
    log(f"  视频时长: {duration//60}分{duration%60}秒")
    
    # 计算分段（每段最多约25分钟音频，64kbps MP3约12MB/25分钟）
    max_chunk_duration = 25 * 60  # 25分钟一段
    
    if duration <= max_chunk_duration:
        # 单段处理
        temp_audio = output_txt.replace(".txt", "_temp.mp3")
        try:
            log(f"  [1/2] 提取音频...")
            extract_audio_segment(video_path, temp_audio)
            audio_size = os.path.getsize(temp_audio)
            log(f"  音频大小: {audio_size/(1024*1024):.1f}MB")
            
            # 如果超过Gemini限制，再分割
            if audio_size > GEMINI_MAX_AUDIO_BYTES:
                # 需要分块，重新用更短的时长提取
                shorter_chunk = 10 * 60  # 10分钟
                chunks_needed = (duration + shorter_chunk - 1) // shorter_chunk
                log(f"  ⚠️ 音频过大，分{chunks_needed}段处理")
                os.remove(temp_audio)
                return process_video_chunked(video_path, output_txt, duration, shorter_chunk, progress)
            
            log(f"  [2/2] Gemini转录...")
            text = gemini_transcribe(temp_audio)
            
            # 清理
            if os.path.exists(temp_audio):
                os.remove(temp_audio)
            
            # 写入结果
            write_transcript(output_txt, video_path, text, duration)
            progress["completed"].append(video_path)
            save_progress(progress)
            return True
            
        except Exception as e:
            log(f"  ❌ 失败: {e}")
            traceback.print_exc()
            progress["failed"].append({"file": video_path, "error": str(e)})
            save_progress(progress)
            if os.path.exists(temp_audio):
                try: os.remove(temp_audio)
                except: pass
            return False
    else:
        # 多段处理
        chunk_duration = max_chunk_duration
        return process_video_chunked(video_path, output_txt, duration, chunk_duration, progress)

def process_video_chunked(video_path, output_txt, duration, chunk_duration, progress):
    """分段处理长视频"""
    chunks = []
    num_chunks = (duration + chunk_duration - 1) // chunk_duration
    
    for i in range(num_chunks):
        start = i * chunk_duration
        dur = min(chunk_duration, duration - start)
        temp_audio = output_txt.replace(".txt", f"_chunk{i+1}_temp.mp3")
        
        try:
            log(f"  段{i+1}/{num_chunks}: {start//60}分{start%60}秒 → {(start+dur)//60}分{(start+dur)%60}秒")
            extract_audio_segment(video_path, temp_audio, start_sec=start, duration_sec=dur)
            
            chunk_info = f"这是音频的第{i+1}段（共{num_chunks}段），时间段约{start//60}分到{(start+dur)//60}分。"
            text = gemini_transcribe(temp_audio, chunk_info)
            chunks.append(text)
            
            if os.path.exists(temp_audio):
                os.remove(temp_audio)
            
            # 段间等待，避免限流
            if i < num_chunks - 1:
                time.sleep(5)
                
        except Exception as e:
            log(f"  ❌ 段{i+1}失败: {e}")
            chunks.append(f"[第{i+1}段转录失败: {e}]")
            if os.path.exists(temp_audio):
                try: os.remove(temp_audio)
                except: pass
    
    full_text = "\n\n".join(chunks)
    write_transcript(output_txt, video_path, full_text, duration, num_chunks)
    progress["completed"].append(video_path)
    save_progress(progress)
    return True

def write_transcript(output_txt, video_path, text, duration, chunks=1):
    """写入转录结果"""
    with open(output_txt, 'w', encoding='utf-8') as f:
        f.write(f"# 转录文本\n")
        f.write(f"# 源文件: {os.path.basename(video_path)}\n")
        f.write(f"# 时长: {duration//60}分{duration%60}秒\n")
        f.write(f"# 转录方式: Gemini API ({GEMINI_MODEL})\n")
        if chunks > 1:
            f.write(f"# 分段数: {chunks}\n")
        f.write(f"# 时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
        f.write(text)

# ===== 文档处理 =====

def extract_pdf_text(pdf_path):
    import pdfplumber
    texts = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                texts.append(f"--- 第{i+1}页 ---\n{text}")
    return "\n\n".join(texts)

def extract_pptx_text(pptx_path):
    from pptx import Presentation
    prs = Presentation(pptx_path)
    texts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_texts.append(row_text)
        if slide_texts:
            texts.append(f"--- 幻灯片 {i+1} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)

def extract_docx_text(docx_path):
    from docx import Document
    doc = Document(docx_path)
    texts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            texts.append(text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                texts.append(row_text)
    return "\n".join(texts)

def extract_xlsx_text(xlsx_path):
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
    texts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_texts = [f"--- 工作表: {sheet_name} ---"]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip(" |"):
                sheet_texts.append(row_text)
        texts.append("\n".join(sheet_texts))
    wb.close()
    return "\n\n".join(texts)

def process_document(doc_path, progress):
    all_dirs = VIDEO_DIRS + DOC_DIRS
    source_root = find_source_root(doc_path, all_dirs)
    output_txt = get_output_path(doc_path, source_root)
    
    if doc_path in progress["completed"]:
        return True
    if os.path.exists(output_txt) and os.path.getsize(output_txt) > 50:
        progress["completed"].append(doc_path)
        save_progress(progress)
        return True
    
    os.makedirs(os.path.dirname(output_txt), exist_ok=True)
    ext = os.path.splitext(doc_path)[1].lower()
    
    try:
        if ext == ".pdf":
            text = extract_pdf_text(doc_path)
        elif ext == ".pptx":
            text = extract_pptx_text(doc_path)
        elif ext == ".docx":
            text = extract_docx_text(doc_path)
        elif ext == ".xlsx":
            text = extract_xlsx_text(doc_path)
        else:
            return True
        
        if not text.strip():
            text = "（提取到的文字为空）"
        
        with open(output_txt, 'w', encoding='utf-8') as f:
            f.write(f"# 文档文字提取\n")
            f.write(f"# 源文件: {os.path.basename(doc_path)}\n")
            f.write(f"# 格式: {ext}\n")
            f.write(f"# 时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            f.write(text)
        
        progress["completed"].append(doc_path)
        save_progress(progress)
        return True
    except Exception as e:
        log(f"  ❌ 文档处理失败: {e}")
        progress["failed"].append({"file": doc_path, "error": str(e)})
        save_progress(progress)
        return False

def scan_files(dirs, extensions):
    files = []
    for d in dirs:
        if not os.path.exists(d):
            log(f"  [警告] 目录不存在: {d}")
            continue
        for root, _, filenames in os.walk(d):
            for fn in filenames:
                if os.path.splitext(fn)[1].lower() in extensions:
                    files.append(os.path.join(root, fn))
    return sorted(files)

def main():
    global GEMINI_API_KEY
    
    print("=" * 70)
    print("  批量转录 v2 - Gemini API + 文档提取")
    print("=" * 70)
    
    # 加载API Key
    GEMINI_API_KEY = load_gemini_key()
    if not GEMINI_API_KEY:
        print("❌ 未找到Gemini API Key")
        sys.exit(1)
    print(f"  Gemini API Key: {GEMINI_API_KEY[:10]}...")
    print(f"  Gemini Model: {GEMINI_MODEL}")
    print(f"  输出目录: {OUTPUT_ROOT}")
    
    # 确保ffmpeg存在
    if not os.path.exists(FFMPEG_PATH):
        print(f"❌ ffmpeg不存在: {FFMPEG_PATH}")
        sys.exit(1)
    
    progress = load_progress()
    
    # 扫描文件
    log("📂 扫描文件...")
    video_files = scan_files(VIDEO_DIRS, {".mp4", ".avi", ".mkv", ".mov", ".flv"})
    doc_files = scan_files(DOC_DIRS, {".pdf", ".pptx", ".docx", ".xlsx"})
    log(f"  视频: {len(video_files)} 个")
    log(f"  文档: {len(doc_files)} 个")
    
    # === 阶段1：文档提取（快） ===
    log(f"\n{'='*60}")
    log(f"  阶段1：文档文字提取 ({len(doc_files)} 个)")
    log(f"{'='*60}")
    
    doc_done = doc_fail = 0
    for i, doc_path in enumerate(doc_files):
        log(f"[{i+1}/{len(doc_files)}] {os.path.basename(doc_path)}")
        if process_document(doc_path, progress):
            doc_done += 1
        else:
            doc_fail += 1
    
    log(f"📊 文档完成: ✅{doc_done} ❌{doc_fail}")
    
    # === 阶段2：视频转录（慢） ===
    log(f"\n{'='*60}")
    log(f"  阶段2：视频转录 ({len(video_files)} 个)")
    log(f"  ⏱️ 使用Gemini API，预计每个视频1-3分钟")
    log(f"{'='*60}")
    
    vid_done = vid_fail = 0
    start_time = time.time()
    
    for i, video_path in enumerate(video_files):
        elapsed = time.time() - start_time
        if vid_done > 0:
            avg = elapsed / vid_done
            remaining = avg * (len(video_files) - i)
            eta = f"ETA ~{remaining/60:.0f}分钟"
        else:
            eta = ""
        
        size_mb = os.path.getsize(video_path) / (1024*1024)
        log(f"\n[{i+1}/{len(video_files)}] {os.path.basename(video_path)} ({size_mb:.0f}MB) {eta}")
        
        if process_video(video_path, progress):
            vid_done += 1
        else:
            vid_fail += 1
        
        # 视频间等待，避免API限流
        time.sleep(3)
    
    total_time = time.time() - start_time
    
    # 汇总
    log(f"\n{'='*60}")
    log(f"  🏁 全部完成！")
    log(f"{'='*60}")
    log(f"  文档: ✅{doc_done} ❌{doc_fail}")
    log(f"  视频: ✅{vid_done} ❌{vid_fail}")
    log(f"  总耗时: {total_time/60:.1f}分钟")
    log(f"  输出: {OUTPUT_ROOT}")
    
    if progress["failed"]:
        log(f"\n❌ 失败列表:")
        for item in progress["failed"]:
            log(f"  - {item['file']}: {item['error']}")

if __name__ == "__main__":
    main()
